"""Routes any file to the correct parser and returns normalized chunks with metadata."""
import hashlib
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches
import openpyxl
import csv


# ── Helpers ──────────────────────────────────────────────────────────────────

def _file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()[:16]


def _base_meta(path: Path, extra: dict | None = None) -> dict[str, Any]:
    """Build the standard metadata block every chunk must carry."""
    stat = path.stat()
    meta = {
        "source_file": path.name,
        "source_path": str(path.absolute()),
        "source_folder": str(path.parent.absolute()),
        "file_type": path.suffix.lower().lstrip("."),
        "file_size_bytes": stat.st_size,
        "file_hash": _file_hash(path),
        "ingested_at": datetime.utcnow().isoformat(),
        "last_modified": datetime.utcfromtimestamp(stat.st_mtime).isoformat(),
    }
    if extra:
        meta.update(extra)
    return meta


def _make_chunk(text: str, meta: dict) -> dict[str, Any]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return None
    return {"text": text, "metadata": meta}


# ── Parsers ───────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> list[dict]:
    chunks = []
    base = _base_meta(path, {"document_type": "pdf"})
    with pdfplumber.open(path) as pdf:
        base["total_pages"] = len(pdf.pages)
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            tables = page.extract_tables()
            if tables:
                for tbl in tables:
                    rows = [" | ".join(str(c) for c in row if c) for row in tbl if any(row)]
                    text += "\n" + "\n".join(rows)
            meta = {**base, "page_number": i, "section": f"Page {i}"}
            chunk = _make_chunk(text, meta)
            if chunk:
                chunks.append(chunk)
    return chunks


def _parse_docx(path: Path) -> list[dict]:
    chunks = []
    base = _base_meta(path, {"document_type": "docx"})
    doc = DocxDocument(str(path))

    current_heading = "Document Start"
    para_buffer = []

    for para in doc.paragraphs:
        style = para.style.name.lower()
        text = para.text.strip()
        if not text:
            continue
        if "heading" in style:
            if para_buffer:
                meta = {**base, "section": current_heading}
                chunk = _make_chunk(" ".join(para_buffer), meta)
                if chunk:
                    chunks.append(chunk)
                para_buffer = []
            current_heading = text
        else:
            para_buffer.append(text)

    if para_buffer:
        meta = {**base, "section": current_heading}
        chunk = _make_chunk(" ".join(para_buffer), meta)
        if chunk:
            chunks.append(chunk)

    # Tables
    for i, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        meta = {**base, "section": f"Table {i}"}
        chunk = _make_chunk("\n".join(rows), meta)
        if chunk:
            chunks.append(chunk)

    return chunks


def _parse_pptx(path: Path) -> list[dict]:
    chunks = []
    base = _base_meta(path, {"document_type": "pptx"})
    prs = Presentation(str(path))
    base["total_slides"] = len(prs.slides)

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    if shape.shape_type == 13:  # title
                        title = t
                    else:
                        texts.append(t)
        if title:
            texts.insert(0, f"Slide Title: {title}")
        meta = {**base, "slide_number": i, "section": title or f"Slide {i}"}
        chunk = _make_chunk(" ".join(texts), meta)
        if chunk:
            chunks.append(chunk)

    return chunks


def _parse_excel(path: Path) -> list[dict]:
    chunks = []
    base = _base_meta(path, {"document_type": path.suffix.lstrip(".")})
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        headers = None
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            row = [str(c) if c is not None else "" for c in row]
            if i == 0:
                headers = row
            else:
                if headers:
                    rows.append(dict(zip(headers, row)))
                else:
                    rows.append(row)
            # chunk every 50 rows to avoid huge chunks
            if len(rows) >= 50:
                text = json.dumps(rows, ensure_ascii=False)
                meta = {**base, "sheet": sheet_name, "section": f"Sheet: {sheet_name}"}
                chunk = _make_chunk(text, meta)
                if chunk:
                    chunks.append(chunk)
                rows = []

        if rows:
            text = json.dumps(rows, ensure_ascii=False)
            meta = {**base, "sheet": sheet_name, "section": f"Sheet: {sheet_name}"}
            chunk = _make_chunk(text, meta)
            if chunk:
                chunks.append(chunk)

    wb.close()
    return chunks


def _parse_csv(path: Path) -> list[dict]:
    chunks = []
    base = _base_meta(path, {"document_type": "csv"})
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.DictReader(f)
        batch = []
        for row in reader:
            batch.append(row)
            if len(batch) >= 50:
                meta = {**base, "section": "CSV Data"}
                chunk = _make_chunk(json.dumps(batch, ensure_ascii=False), meta)
                if chunk:
                    chunks.append(chunk)
                batch = []
        if batch:
            meta = {**base, "section": "CSV Data"}
            chunk = _make_chunk(json.dumps(batch, ensure_ascii=False), meta)
            if chunk:
                chunks.append(chunk)
    return chunks


def _parse_text(path: Path) -> list[dict]:
    base = _base_meta(path, {"document_type": path.suffix.lstrip(".") or "txt"})
    text = path.read_text(encoding="utf-8", errors="replace")
    chunk = _make_chunk(text, {**base, "section": "Full Document"})
    return [chunk] if chunk else []


def _parse_json(path: Path) -> list[dict]:
    base = _base_meta(path, {"document_type": "json"})
    data = json.loads(path.read_text(encoding="utf-8"))
    text = json.dumps(data, indent=2, ensure_ascii=False)
    chunk = _make_chunk(text, {**base, "section": "JSON Document"})
    return [chunk] if chunk else []


# ── Router ────────────────────────────────────────────────────────────────────

PARSER_MAP = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".doc":  _parse_docx,
    ".pptx": _parse_pptx,
    ".ppt":  _parse_pptx,
    ".xlsx": _parse_excel,
    ".xls":  _parse_excel,
    ".csv":  _parse_csv,
    ".txt":  _parse_text,
    ".md":   _parse_text,
    ".rtf":  _parse_text,
    ".html": _parse_text,
    ".htm":  _parse_text,
    ".json": _parse_json,
}


def parse_document(path: str | Path) -> list[dict[str, Any]]:
    """
    Parse any supported document.
    Returns a list of dicts: {"text": str, "metadata": dict}
    Every chunk always carries full source file/path/folder info.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    parser = PARSER_MAP.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")

    try:
        return parser(path)
    except Exception as e:
        raise RuntimeError(f"Failed to parse {path.name}: {e}") from e


def get_supported_extensions() -> set[str]:
    return set(PARSER_MAP.keys())
