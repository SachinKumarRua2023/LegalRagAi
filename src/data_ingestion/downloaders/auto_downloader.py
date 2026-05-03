"""
Auto-downloader for US legal data.
Sources (all free, no API key needed):
  1. Oyez.org API  — Supreme Court cases (JSON)
  2. HuggingFace   — casehold dataset (legal holdings, Parquet)
Saves as: TXT, JSON, DOCX, PPTX, XLSX, CSV
"""
from __future__ import annotations
import json, csv, time, re
from pathlib import Path
from typing import Any

import requests
from rich.console import Console

from config.settings import DATA_RAW_PATH

console = Console()


# ── helpers ───────────────────────────────────────────────────────────────────

def _clean(text: str) -> str:
    text = re.sub(r"<[^>]+>", " ", text or "")
    return re.sub(r"\s+", " ", text).strip()


# ── Source 1 : Oyez.org — SCOTUS cases ───────────────────────────────────────

OYEZ_LIST  = "https://api.oyez.org/cases"
OYEZ_CASE  = "https://api.oyez.org/cases/{term}/{docket}"
OYEZ_HDR   = {"User-Agent": "CompleteRagAI/1.0 (legal-research; contact: research@example.com)"}


def _oyez_fetch(max_cases: int = 200) -> list[dict]:
    """Fetch SCOTUS case data from Oyez public API."""
    console.print("  [cyan]Fetching SCOTUS cases from Oyez.org...[/cyan]")
    records, page = [], 0

    while len(records) < max_cases:
        try:
            r = requests.get(
                OYEZ_LIST,
                params={"per_page": 50, "page": page, "filter": ""},
                headers=OYEZ_HDR,
                timeout=20,
            )
            r.raise_for_status()
            cases = r.json()
        except Exception as e:
            console.print(f"  [yellow]Oyez list error: {e}[/yellow]")
            break

        if not cases:
            break

        for case in cases:
            if len(records) >= max_cases:
                break
            rec = {
                "case_name": case.get("name", ""),
                "term": str(case.get("term", "")),
                "docket": str(case.get("docket_number", "")),
                "topic": "scotus",
                "source": "oyez.org",
                "description": case.get("description", ""),
                "facts": _clean(case.get("facts_of_the_case", "") or ""),
                "question": _clean(case.get("question", "") or ""),
                "conclusion": _clean(case.get("conclusion", "") or ""),
                "url": case.get("href", ""),
            }
            # Combine into full text
            parts = []
            if rec["case_name"]:  parts.append(f"CASE: {rec['case_name']}")
            if rec["term"]:       parts.append(f"TERM: {rec['term']}")
            if rec["docket"]:     parts.append(f"DOCKET: {rec['docket']}")
            if rec["description"]:parts.append(f"\nDESCRIPTION:\n{rec['description']}")
            if rec["facts"]:      parts.append(f"\nFACTS:\n{rec['facts']}")
            if rec["question"]:   parts.append(f"\nQUESTION:\n{rec['question']}")
            if rec["conclusion"]: parts.append(f"\nCONCLUSION:\n{rec['conclusion']}")
            rec["text"] = "\n".join(parts)

            if len(rec["text"]) > 80:
                records.append(rec)

        page += 1
        time.sleep(0.2)

    console.print(f"  [green]Oyez: {len(records)} SCOTUS cases fetched[/green]")
    return records


# ── Source 2 : HuggingFace — casehold (Parquet, no script) ───────────────────

def _hf_casehold(max_records: int = 200) -> list[dict]:
    """Download CaseHOLD legal holdings dataset from HuggingFace."""
    try:
        from datasets import load_dataset
    except ImportError:
        return []

    console.print("  [cyan]Fetching CaseHOLD from HuggingFace...[/cyan]")
    try:
        ds = load_dataset("casehold/casehold", "all", split="train", streaming=True)
    except Exception as e:
        console.print(f"  [yellow]CaseHOLD failed: {e}[/yellow]")
        return []

    records = []
    for row in ds:
        if len(records) >= max_records:
            break
        citing = _clean(str(row.get("citing_prompt", "")))
        if len(citing) < 80:
            continue
        holdings = [_clean(str(row.get(f"holding_{i}", ""))) for i in range(5)]
        holdings = [h for h in holdings if h]

        rec = {
            "text": f"LEGAL HOLDING CONTEXT:\n{citing}\n\nHOLDINGS:\n" + "\n".join(f"{i+1}. {h}" for i, h in enumerate(holdings)),
            "topic": "casehold",
            "source": "casehold/casehold",
            "case_name": f"CaseHOLD Record {len(records)+1}",
            "description": "Legal case holding from CaseHOLD dataset",
        }
        records.append(rec)

    console.print(f"  [green]CaseHOLD: {len(records)} legal holdings fetched[/green]")
    return records


# ── File format converters ────────────────────────────────────────────────────

def _to_txt(records: list[dict], out: Path, slug: str) -> list[Path]:
    saved = []
    for i, r in enumerate(records):
        p = out / f"{slug}_{i+1:04d}.txt"
        header = (f"SOURCE: {r.get('source','')}\n"
                  f"CASE: {r.get('case_name','')}\n"
                  f"TOPIC: {r.get('topic','')}\n"
                  + ("-" * 60) + "\n\n")
        p.write_text(header + r["text"], encoding="utf-8")
        saved.append(p)
    return saved


def _to_json(records: list[dict], out: Path, slug: str) -> list[Path]:
    p = out / f"{slug}_cases.json"
    p.write_text(json.dumps(records, indent=2, ensure_ascii=False), encoding="utf-8")
    return [p]


def _to_docx(records: list[dict], out: Path, slug: str) -> list[Path]:
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
    except ImportError:
        return []

    doc = Document()
    doc.add_heading(f"US Legal Cases - {slug.replace('_',' ').title()}", 0)

    for i, r in enumerate(records[:80], 1):
        doc.add_heading(f"{i}. {r.get('case_name', f'Case {i}')}", level=1)
        meta = doc.add_paragraph()
        run = meta.add_run(f"Source: {r.get('source','')}  |  Topic: {r.get('topic','')}")
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        doc.add_paragraph(r["text"][:1000] + ("..." if len(r["text"]) > 1000 else ""))
        doc.add_paragraph()

    p = out / f"{slug}_cases.docx"
    doc.save(str(p))
    return [p]


def _to_pptx(records: list[dict], out: Path, slug: str) -> list[Path]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return []

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = f"Legal Cases: {slug.replace('_',' ').title()}"
    sl.placeholders[1].text = f"{len(records)} cases - CompleteRagAI"

    for i, r in enumerate(records[:20], 1):
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        tb = sl.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(1.0))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = r.get("case_name", f"Case {i}")[:80]
        run.font.size = Pt(20); run.font.bold = True

        tb2 = sl.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.7))
        tb2.text_frame.word_wrap = True
        tb2.text_frame.text = r["text"][:700]

    p = out / f"{slug}_cases.pptx"
    prs.save(str(p))
    return [p]


def _to_xlsx(all_records: dict[str, list[dict]], out: Path) -> list[Path]:
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
    except ImportError:
        return []

    wb = openpyxl.Workbook()
    ws_all = wb.active
    ws_all.title = "All Cases"
    hdrs = ["#", "Case Name", "Topic", "Source", "Text Preview"]
    ws_all.append(hdrs)
    hf = Font(bold=True, color="FFFFFF")
    hfill = PatternFill("solid", fgColor="1E3A5F")
    for cell in ws_all[1]:
        cell.font = hf; cell.fill = hfill
        cell.alignment = Alignment(horizontal="center")

    n = 1
    for slug, records in all_records.items():
        ws = wb.create_sheet(slug[:28])
        ws.append(hdrs)
        for cell in ws[1]:
            cell.font = hf; cell.fill = hfill
        for r in records:
            row = [n, r.get("case_name",""), r.get("topic",""), r.get("source",""), r["text"][:400]]
            ws.append(row); ws_all.append(row); n += 1
        ws.column_dimensions["E"].width = 80

    ws_all.column_dimensions["E"].width = 80
    p = out / "all_legal_cases.xlsx"
    wb.save(str(p))
    return [p]


def _to_csv(all_records: dict[str, list[dict]], out: Path) -> list[Path]:
    p = out / "all_legal_cases.csv"
    with open(p, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=["id","case_name","topic","source","text"])
        w.writeheader()
        i = 1
        for records in all_records.values():
            for r in records:
                w.writerow({"id":i, "case_name":r.get("case_name",""),
                            "topic":r.get("topic",""), "source":r.get("source",""),
                            "text":r["text"][:600]})
                i += 1
    return [p]


def _to_pdf(records: list[dict], out: Path, slug: str) -> list[Path]:
    """Create PDF document with legal cases."""
    try:
        from fpdf import FPDF
    except ImportError:
        return _to_pdf_reportlab(records, out, slug)

    def _sanitize_for_pdf(text: str) -> str:
        """Remove/replace Unicode characters that PDF fonts can't handle."""
        if not text:
            return ""
        # Replace common Unicode quotes with ASCII equivalents
        replacements = {
            ''': "'", ''': "'", '"': '"', '"': '"',
            '—': '-', '–': '-', '…': '...',
            '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"',
            '\xa0': ' ',  # Non-breaking space
        }
        for old, new in replacements.items():
            text = text.replace(old, new)
        # Remove any remaining non-ASCII characters
        return text.encode('ascii', 'ignore').decode('ascii')

    class LegalPDF(FPDF):
        def header(self):
            self.set_font('helvetica', 'B', 12)
            title = _sanitize_for_pdf(f'US Legal Cases - {slug.replace("_", " ").title()}')
            self.cell(0, 10, title, 0, 1, 'C')
            self.ln(5)

        def footer(self):
            self.set_y(-15)
            self.set_font('helvetica', 'I', 8)
            self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

    pdf = LegalPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    for i, r in enumerate(records[:100], 1):  # Limit to 100 cases for PDF size
        pdf.set_font('helvetica', 'B', 11)
        case_title = _sanitize_for_pdf(f"{i}. {r.get('case_name', f'Case {i}')[:80]}")
        pdf.cell(0, 8, case_title, 0, 1)

        pdf.set_font('helvetica', '', 9)
        pdf.set_text_color(100, 100, 100)
        meta = _sanitize_for_pdf(f"Source: {r.get('source','')} | Topic: {r.get('topic','')}")
        pdf.cell(0, 5, meta, 0, 1)
        pdf.set_text_color(0, 0, 0)

        pdf.set_font('helvetica', '', 10)
        text = r["text"][:2000] if len(r["text"]) > 2000 else r["text"]
        text = _sanitize_for_pdf(text)
        pdf.multi_cell(0, 5, text)
        pdf.ln(5)

    p = out / f"{slug}_cases.pdf"
    pdf.output(str(p))
    return [p]


def _to_pdf_reportlab(records: list[dict], out: Path, slug: str) -> list[Path]:
    """Fallback PDF generation using reportlab."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError:
        return []

    p = out / f"{slug}_cases.pdf"
    doc = SimpleDocTemplate(str(p), pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    title = Paragraph(f"<b>US Legal Cases - {slug.replace('_', ' ').title()}</b>", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 12))

    for i, r in enumerate(records[:50], 1):
        case_title = f"<b>{i}. {r.get('case_name', f'Case {i}')}</b>"
        story.append(Paragraph(case_title, styles['Heading2']))

        meta = f"<i>Source: {r.get('source','')} | Topic: {r.get('topic','')}</i>"
        story.append(Paragraph(meta, styles['Normal']))

        text = r["text"][:1500].replace('\n', '<br/>')
        story.append(Paragraph(text, styles['Normal']))
        story.append(Spacer(1, 12))

    doc.build(story)
    return [p]


# ── Main ──────────────────────────────────────────────────────────────────────

def download_mixed_legal_data(
    cases_per_topic: int = 200,
    output_dir: Path | None = None,
) -> dict[str, Any]:
    """
    Download US legal data from Oyez + HuggingFace (free, no API keys).
    Saves as TXT, JSON, DOCX, PPTX, XLSX, CSV.
    """
    output_dir = output_dir or (DATA_RAW_PATH / "mixed")
    output_dir.mkdir(parents=True, exist_ok=True)

    all_records: dict[str, list[dict]] = {}
    saved_files: list[str] = []

    console.print(f"\n[bold cyan]Downloading US legal data -> {output_dir}[/bold cyan]")

    # ── SCOTUS from Oyez ──
    oyez_dir = output_dir / "scotus"
    oyez_dir.mkdir(exist_ok=True)
    scotus = _oyez_fetch(max_cases=min(cases_per_topic, 200))
    if scotus:
        all_records["scotus"] = scotus
        console.print("  [dim]Saving SCOTUS as TXT / JSON / DOCX / PPTX / PDF...[/dim]")
        for fp in _to_txt(scotus, oyez_dir, "scotus"):   saved_files.append(str(fp))
        for fp in _to_json(scotus, oyez_dir, "scotus"):  saved_files.append(str(fp))
        for fp in _to_docx(scotus, oyez_dir, "scotus"):  saved_files.append(str(fp))
        for fp in _to_pptx(scotus, oyez_dir, "scotus"):  saved_files.append(str(fp))
        for fp in _to_pdf(scotus, oyez_dir, "scotus"):   saved_files.append(str(fp))

    # ── CaseHOLD from HuggingFace ──
    ch_dir = output_dir / "casehold"
    ch_dir.mkdir(exist_ok=True)
    casehold = _hf_casehold(max_records=min(cases_per_topic, 200))
    if casehold:
        all_records["casehold"] = casehold
        console.print("  [dim]Saving CaseHOLD as TXT / JSON / DOCX / PPTX / PDF...[/dim]")
        for fp in _to_txt(casehold, ch_dir, "casehold"):  saved_files.append(str(fp))
        for fp in _to_json(casehold, ch_dir, "casehold"): saved_files.append(str(fp))
        for fp in _to_docx(casehold, ch_dir, "casehold"): saved_files.append(str(fp))
        for fp in _to_pptx(casehold, ch_dir, "casehold"): saved_files.append(str(fp))
        for fp in _to_pdf(casehold, ch_dir, "casehold"):  saved_files.append(str(fp))

    # ── Cross-topic XLSX + CSV ──
    if all_records:
        console.print("  [dim]Building XLSX + CSV summary...[/dim]")
        for fp in _to_xlsx(all_records, output_dir): saved_files.append(str(fp))
        for fp in _to_csv(all_records, output_dir):  saved_files.append(str(fp))

    total = sum(len(v) for v in all_records.values())
    console.print(f"\n[bold green]Done! {total} cases, {len(saved_files)} files in {output_dir}[/bold green]")

    return {
        "status": "ok",
        "output_dir": str(output_dir),
        "topics": list(all_records.keys()),
        "total_cases": total,
        "files_created": len(saved_files),
        "file_types": ["txt","json","docx","pptx","pdf","xlsx","csv"],
    }
